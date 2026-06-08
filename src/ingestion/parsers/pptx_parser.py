"""
PowerPoint (.pptx) 文档解析器。

使用 python-pptx 提取幻灯片文本。
每张幻灯片生成一个 LlamaDocument，包含标题和正文。
"""

import logging
from typing import List

from llama_index.core.schema import Document as LlamaDocument

from src.ingestion.parsers.base_parser import BaseParser, make_doc_id

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """PowerPoint 解析器——基于 python-pptx，每张幻灯片一个文档。

    Usage::

        parser = PptxParser()
        docs = parser.parse("/path/to/slides.pptx")
        for d in docs:
            print(f"Slide {d.metadata['slide_number']}: {d.metadata['slide_title']}")
    """

    @property
    def supported_format(self) -> str:
        return "pptx"

    def parse(self, source: str) -> List[LlamaDocument]:
        """解析 .pptx 文件。

        Args:
            source: .pptx 文件路径。

        Returns:
            LlamaDocument 列表，每张幻灯片一个文档。
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx 未安装")
            return []

        doc_id = make_doc_id(source)
        doc_name = self._source_name(source)
        documents: List[LlamaDocument] = []

        try:
            prs = Presentation(source)
        except Exception as e:
            logger.error("无法打开 PPT 文件 %s: %s", source, e)
            return []

        total_slides = len(prs.slides)
        total_chars = 0

        for slide_idx, slide in enumerate(prs.slides):
            # 提取所有文本框
            texts: List[str] = []
            slide_title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            # 第一个非空文本作为标题
                            if not slide_title and shape.has_text_frame:
                                slide_title = para_text
                            texts.append(para_text)

            if not texts:
                logger.debug("第 %d 张幻灯片无文本，跳过", slide_idx + 1)
                continue

            full_text = "\n".join(texts)
            total_chars += len(full_text)

            doc = LlamaDocument(
                text=full_text,
                metadata={
                    "doc_id": doc_id,
                    "doc_name": doc_name,
                    "source": source,
                    "format": "pptx",
                    "slide_number": slide_idx + 1,  # 1-based
                    "total_slides": total_slides,
                    "slide_title": slide_title,
                },
            )
            documents.append(doc)

        logger.info(
            "PPT 解析完成: %s (%d/%d 张非空幻灯片, %d 字符)",
            doc_name, len(documents), total_slides, total_chars,
        )
        return documents

    # ------------------------------------------------------------------
    # LlamaIndex 回退
    # ------------------------------------------------------------------

    @staticmethod
    def parse_via_llamaindex(source: str) -> List[LlamaDocument]:
        """使用 LlamaIndex PptxReader 作为备选方案。

        Args:
            source: 文件路径。

        Returns:
            LlamaDocument 列表。
        """
        try:
            from llama_index.readers.file import PptxReader

            reader = PptxReader()
            docs = reader.load_data(file_path=source)  # type: ignore[arg-type]
            logger.info("LlamaIndex PptxReader 解析完成: %s", source)
            return list(docs)
        except ImportError:
            logger.error("llama-index-readers-file 未安装")
            return []
